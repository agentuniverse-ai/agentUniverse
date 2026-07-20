# !/usr/bin/env python3
# -*- coding:utf-8 -*-

# @Time    : 2025/10/28 21:21
# @Author  : Saladday
# @Email   : fanjing.luo@zju.edu.cn
# @FileName: test_zip_reader.py
import io
import tempfile
import unittest
import zipfile
from pathlib import Path

from agentuniverse.agent.action.knowledge.reader.file.zip_reader import ZipReader


class TestZipReader(unittest.TestCase):
    def setUp(self) -> None:
        self.temp_dir = tempfile.TemporaryDirectory()
        self.reader = ZipReader()

    def tearDown(self) -> None:
        self.temp_dir.cleanup()

    def _create_docx_file(self, text: str) -> bytes:
        try:
            from docx import Document
            doc = Document()
            doc.add_paragraph(text)
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.read()
        except ImportError:
            return b""

    def _create_pdf_file(self, text: str) -> bytes:
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            buffer = io.BytesIO()
            pdf = canvas.Canvas(buffer, pagesize=letter)
            pdf.drawString(100, 750, text)
            pdf.save()
            buffer.seek(0)
            return buffer.read()
        except ImportError:
            return b""

    def _create_pptx_file(self, text: str) -> bytes:
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = text
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.read()
        except ImportError:
            return b""

    def _create_xlsx_file(self) -> bytes:
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws['A1'] = '姓名'
            ws['B1'] = '年龄'
            ws['A2'] = '张三'
            ws['B2'] = 25
            ws['A3'] = '李四'
            ws['B3'] = 30
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            return buffer.read()
        except ImportError:
            return b""

    def test_complex_nested_zip_structure(self) -> None:
        archive_path = Path(self.temp_dir.name) / "complex_archive.zip"
        
        level3_zip = io.BytesIO()
        with zipfile.ZipFile(level3_zip, "w") as z3:
            z3.writestr("deep/secret.txt", "这是第三层深度的秘密文档")
            z3.writestr("deep/config.json", '{"level": 3, "type": "configuration"}')
        
        level2_zip = io.BytesIO()
        with zipfile.ZipFile(level2_zip, "w") as z2:
            z2.writestr("reports/report.md", "# 第二层报告\n\n这是嵌套的markdown文档")
            z2.writestr("data/metrics.txt", "CPU: 85%\nMemory: 60%\nDisk: 40%")
            z2.writestr("archives/level3.zip", level3_zip.getvalue())
        
        with zipfile.ZipFile(archive_path, "w") as main_zip:
            main_zip.writestr("README.md", "# 主文档\n\n这是根目录的说明文件")
            main_zip.writestr("docs/intro.txt", "欢迎使用复杂压缩包测试系统")
            main_zip.writestr("docs/guide.md", "## 使用指南\n\n1. 解压文件\n2. 阅读文档\n3. 运行测试")
            
            main_zip.writestr("src/main.py", "def main():\n    print('Hello from ZIP')\n\nif __name__ == '__main__':\n    main()")
            main_zip.writestr("src/utils.py", "def helper():\n    return 'utility function'")
            
            main_zip.writestr("config/settings.json", '{"app": "test", "version": "1.0.0"}')
            main_zip.writestr("config/database.yml", "host: localhost\nport: 5432\ndatabase: testdb")
            
            main_zip.writestr("data/sample.csv", "Name,Age,City\nAlice,28,Beijing\nBob,32,Shanghai\nCarol,25,Guangzhou")
            
            docx_content = self._create_docx_file("这是一个Word文档，包含重要信息")
            if docx_content:
                main_zip.writestr("documents/report.docx", docx_content)
            
            pdf_content = self._create_pdf_file("这是PDF文档的内容")
            if pdf_content:
                main_zip.writestr("documents/presentation.pdf", pdf_content)
            
            pptx_content = self._create_pptx_file("项目演示PPT")
            if pptx_content:
                main_zip.writestr("documents/slides.pptx", pptx_content)
            
            xlsx_content = self._create_xlsx_file()
            if xlsx_content:
                main_zip.writestr("data/employees.xlsx", xlsx_content)
            
            main_zip.writestr("logs/app.log", "[INFO] Application started\n[DEBUG] Loading configuration\n[INFO] Ready")
            main_zip.writestr("logs/error.log", "[ERROR] Sample error message")
            
            main_zip.writestr("web/index.html", "<html><body><h1>欢迎</h1></body></html>")
            main_zip.writestr("web/style.css", "body { font-family: Arial; }")
            
            main_zip.writestr("nested_archives/level2.zip", level2_zip.getvalue())
        
        docs = self.reader._load_data(archive_path)
        
        self.assertGreater(len(docs), 0)
        
        file_names = [doc.metadata.get("file_name") for doc in docs]
        archive_paths = [doc.metadata.get("archive_path") for doc in docs]
        
        self.assertIn("README.md", file_names)
        self.assertIn("main.py", file_names)
        self.assertIn("settings.json", file_names)
        
        nested_docs = [d for d in docs if "level2.zip" in d.metadata.get("archive_path", "")]
        self.assertGreater(len(nested_docs), 0)
        
        deep_nested = [d for d in docs if "level3.zip" in d.metadata.get("archive_path", "")]
        self.assertGreater(len(deep_nested), 0)
        
        txt_docs = [d for d in docs if d.metadata.get("file_name", "").endswith(".txt")]
        self.assertGreater(len(txt_docs), 0)
        
        py_docs = [d for d in docs if d.metadata.get("file_name", "").endswith(".py")]
        self.assertEqual(len(py_docs), 2)
        
        depths = [doc.metadata.get("archive_depth", 0) for doc in docs]
        self.assertIn(0, depths)
        self.assertIn(1, depths)
        self.assertIn(2, depths)

    def test_load_text_file(self) -> None:
        archive_path = Path(self.temp_dir.name) / "sample.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("docs/readme.txt", "hello world")
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 1)
        doc = docs[0]
        self.assertEqual(doc.text, "hello world")
        self.assertEqual(doc.metadata["file_name"], "readme.txt")
        self.assertEqual(doc.metadata["archive_root"], "sample.zip")
        self.assertEqual(doc.metadata["archive_path"], "docs/readme.txt")

    def test_nested_zip(self) -> None:
        archive_path = Path(self.temp_dir.name) / "nested.zip"
        nested_buffer = io.BytesIO()
        with zipfile.ZipFile(nested_buffer, "w") as nested:
            nested.writestr("inner/data.txt", "nested data")
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("folder/archive.zip", nested_buffer.getvalue())
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 1)
        doc = docs[0]
        self.assertEqual(doc.text, "nested data")
        self.assertEqual(doc.metadata["archive_path"], "folder/archive.zip/inner/data.txt")

    def test_multiple_file_types(self) -> None:
        archive_path = Path(self.temp_dir.name) / "mixed.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("document.txt", "文本内容")
            archive.writestr("readme.md", "# Markdown标题\n正文内容")
            archive.writestr("code.py", "print('Python代码')")
            archive.writestr("data.json", '{"key": "value"}')
            archive.writestr("config.yml", "setting: true")
            archive.writestr("data.csv", "列1,列2\n值1,值2")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 6)
        
        extensions = {doc.metadata["file_name"].split(".")[-1] for doc in docs}
        self.assertIn("txt", extensions)
        self.assertIn("md", extensions)
        self.assertIn("py", extensions)
        self.assertIn("json", extensions)

    def test_exceeds_file_size_limit(self) -> None:
        archive_path = Path(self.temp_dir.name) / "limit.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("large.txt", "a" * 4096)
        limited_reader = ZipReader(max_file_size=1024, max_total_size=2048)
        with self.assertRaises(ValueError):
            limited_reader._load_data(archive_path)

    def test_exceeds_depth_limit(self) -> None:
        archive_path = Path(self.temp_dir.name) / "deep.zip"
        
        current = io.BytesIO()
        with zipfile.ZipFile(current, "w") as z:
            z.writestr("data.txt", "deepest")
        
        for i in range(10):
            parent = io.BytesIO()
            with zipfile.ZipFile(parent, "w") as z:
                z.writestr(f"level{i}.zip", current.getvalue())
            current = parent
        
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("nested.zip", current.getvalue())
        
        shallow_reader = ZipReader(max_depth=2)
        with self.assertRaises(ValueError):
            shallow_reader._load_data(archive_path)

    def test_compression_ratio_limit(self) -> None:
        archive_path = Path(self.temp_dir.name) / "compressed.zip"
        highly_compressible = "a" * 100000
        
        with zipfile.ZipFile(archive_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("repetitive.txt", highly_compressible)
        
        strict_reader = ZipReader(max_compression_ratio=10)
        with self.assertRaises(ValueError):
            strict_reader._load_data(archive_path)

    def test_custom_metadata(self) -> None:
        archive_path = Path(self.temp_dir.name) / "meta.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("file.txt", "content")
        
        custom_meta = {
            "source": "测试来源",
            "category": "文档类别",
            "priority": "高"
        }
        
        docs = self.reader._load_data(archive_path, ext_info=custom_meta)
        self.assertEqual(len(docs), 1)
        doc = docs[0]
        self.assertEqual(doc.metadata["source"], "测试来源")
        self.assertEqual(doc.metadata["category"], "文档类别")
        self.assertEqual(doc.metadata["priority"], "高")

    def test_empty_files_ignored(self) -> None:
        archive_path = Path(self.temp_dir.name) / "empty.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("empty.txt", "")
            archive.writestr("not_empty.txt", "有内容")
        
        docs = self.reader._load_data(archive_path)
        self.assertGreater(len(docs), 0)
        non_empty_docs = [d for d in docs if d.text.strip()]
        self.assertEqual(len(non_empty_docs), 1)
        self.assertEqual(non_empty_docs[0].text, "有内容")

    def test_special_characters_in_path(self) -> None:
        archive_path = Path(self.temp_dir.name) / "special.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("中文目录/文件名.txt", "中文内容")
            archive.writestr("folder with spaces/file name.txt", "content")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 2)
        file_names = [doc.metadata["file_name"] for doc in docs]
        self.assertIn("文件名.txt", file_names)
        self.assertIn("file name.txt", file_names)

    def test_ultra_complex_nested_structure(self) -> None:
        archive_path = Path(self.temp_dir.name) / "ultra_complex.zip"
        
        level4_zip = io.BytesIO()
        with zipfile.ZipFile(level4_zip, "w") as z4:
            z4.writestr("final/ultimate.txt", "最深层文档内容")
            z4.writestr("final/data.json", '{"depth": 4}')
            z4.writestr("final/script.py", "print('level 4')")
        
        level3_zip = io.BytesIO()
        with zipfile.ZipFile(level3_zip, "w") as z3:
            z3.writestr("deep/secret.txt", "第三层秘密")
            z3.writestr("deep/config.yml", "level: 3\ntype: config")
            z3.writestr("deep/code.py", "def level3(): pass")
            z3.writestr("archives/level4.zip", level4_zip.getvalue())
        
        level2_zip = io.BytesIO()
        with zipfile.ZipFile(level2_zip, "w") as z2:
            z2.writestr("reports/report.md", "# 第二层报告")
            z2.writestr("reports/summary.txt", "总结内容")
            z2.writestr("data/metrics.csv", "Name,Value\nCPU,85\nMemory,60")
            z2.writestr("data/analysis.json", '{"status": "ok"}')
            z2.writestr("scripts/process.py", "def process(): return True")
            z2.writestr("archives/level3.zip", level3_zip.getvalue())
        
        level1_zip = io.BytesIO()
        with zipfile.ZipFile(level1_zip, "w") as z1:
            z1.writestr("docs/readme.md", "# Level 1 文档")
            z1.writestr("docs/notes.txt", "笔记内容")
            z1.writestr("code/main.py", "def main(): print('level1')")
            z1.writestr("nested/level2.zip", level2_zip.getvalue())
        
        with zipfile.ZipFile(archive_path, "w") as main_zip:
            main_zip.writestr("README.md", "# 超级复杂压缩包\n\n包含4层嵌套结构")
            main_zip.writestr("LICENSE.txt", "MIT License")
            main_zip.writestr("docs/intro.md", "## 介绍\n\n这是一个复杂的测试")
            main_zip.writestr("docs/guide.md", "## 指南\n\n使用说明")
            main_zip.writestr("src/app.py", "class App:\n    def run(self): pass")
            main_zip.writestr("src/utils.py", "def helper(): return 42")
            main_zip.writestr("src/config.py", "CONFIG = {'key': 'value'}")
            main_zip.writestr("config/app.json", '{"name": "test"}')
            main_zip.writestr("config/db.yml", "database: test")
            main_zip.writestr("data/input.csv", "A,B,C\n1,2,3\n4,5,6")
            main_zip.writestr("data/output.txt", "结果数据")
            main_zip.writestr("tests/test_app.py", "def test_run(): assert True")
            main_zip.writestr("archives/level1.zip", level1_zip.getvalue())
        
        docs = self.reader._load_data(archive_path)
        
        self.assertGreater(len(docs), 20)
        
        depths = [doc.metadata.get("archive_depth", 0) for doc in docs]
        self.assertIn(0, depths)
        self.assertIn(1, depths)
        self.assertIn(2, depths)
        self.assertIn(3, depths)
        
        level4_docs = [d for d in docs if "level4.zip" in d.metadata.get("archive_path", "")]
        self.assertGreater(len(level4_docs), 0)
        
        py_files = [d for d in docs if d.metadata.get("file_name", "").endswith(".py")]
        self.assertGreater(len(py_files), 5)
        
        md_files = [d for d in docs if d.metadata.get("file_name", "").endswith(".md")]
        self.assertGreater(len(md_files), 3)

    def test_code_files_extraction(self) -> None:
        archive_path = Path(self.temp_dir.name) / "code_archive.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("main.py", "#!/usr/bin/env python3\nprint('Python')")
            archive.writestr("app.js", "console.log('JavaScript');")
            archive.writestr("types.ts", "interface User { name: string; }")
            archive.writestr("Main.java", "public class Main { }")
            archive.writestr("main.go", "package main\nfunc main() {}")
            archive.writestr("lib.cpp", "#include <iostream>\nint main() {}")
            archive.writestr("utils.rs", "fn main() { println!(\"Rust\"); }")
            archive.writestr("script.sh", "#!/bin/bash\necho 'Shell'")
            archive.writestr("app.rb", "puts 'Ruby'")
            archive.writestr("index.php", "<?php echo 'PHP'; ?>")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 10)
        
        languages = [doc.metadata.get("language") for doc in docs]
        self.assertIn("python", languages)
        self.assertIn("javascript", languages)
        self.assertIn("typescript", languages)
        self.assertIn("java", languages)
        self.assertIn("go", languages)
        self.assertIn("cpp", languages)
        self.assertIn("rust", languages)
        self.assertIn("shell", languages)

    def test_mixed_documents_extraction(self) -> None:
        archive_path = Path(self.temp_dir.name) / "docs_archive.zip"
        
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("notes.txt", "这是文本笔记\n第二行内容")
            archive.writestr("readme.md", "# 项目说明\n\n## 功能\n- 功能1\n- 功能2")
            archive.writestr("data.csv", "姓名,年龄,城市\n张三,25,北京\n李四,30,上海\n王五,28,深圳")
            archive.writestr("config.json", '{\n  "version": "1.0",\n  "author": "测试"\n}')
            archive.writestr("settings.yml", "debug: true\nport: 8080\nhost: localhost")
            archive.writestr("index.html", "<html><body><h1>标题</h1><p>段落</p></body></html>")
            archive.writestr("style.css", "body { margin: 0; padding: 0; }\nh1 { color: blue; }")
            archive.writestr("data.xml", '<?xml version="1.0"?>\n<root><item>数据</item></root>')
            archive.writestr("app.log", "[2025-10-28 10:00:00] INFO: 应用启动\n[2025-10-28 10:00:01] DEBUG: 初始化完成")
            
            docx_content = self._create_docx_file("Word文档测试内容\n包含多行文字")
            if docx_content:
                archive.writestr("report.docx", docx_content)
            
            pdf_content = self._create_pdf_file("PDF测试文档")
            if pdf_content:
                archive.writestr("document.pdf", pdf_content)
        
        docs = self.reader._load_data(archive_path)
        self.assertGreater(len(docs), 8)
        
        file_types = {doc.metadata.get("file_name", "").split(".")[-1] for doc in docs}
        self.assertIn("txt", file_types)
        self.assertIn("md", file_types)
        self.assertIn("csv", file_types)
        self.assertIn("json", file_types)

    def test_deeply_nested_directories(self) -> None:
        archive_path = Path(self.temp_dir.name) / "deep_dirs.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("level1/file1.txt", "内容1")
            archive.writestr("level1/level2/file2.txt", "内容2")
            archive.writestr("level1/level2/level3/file3.txt", "内容3")
            archive.writestr("level1/level2/level3/level4/file4.txt", "内容4")
            archive.writestr("level1/level2/level3/level4/level5/file5.txt", "内容5")
            archive.writestr("a/b/c/d/e/f/g/deep.txt", "很深的目录")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 6)
        
        paths = [doc.metadata.get("archive_path", "") for doc in docs]
        self.assertTrue(any("level5" in p for p in paths))
        self.assertTrue(any("a/b/c/d/e/f/g" in p for p in paths))

    def test_duplicate_filenames_different_paths(self) -> None:
        archive_path = Path(self.temp_dir.name) / "duplicates.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("dir1/config.txt", "配置1")
            archive.writestr("dir2/config.txt", "配置2")
            archive.writestr("dir3/config.txt", "配置3")
            archive.writestr("a/b/readme.md", "# 说明A")
            archive.writestr("c/d/readme.md", "# 说明B")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 5)
        
        config_docs = [d for d in docs if d.metadata.get("file_name") == "config.txt"]
        self.assertEqual(len(config_docs), 3)
        
        paths = [d.metadata.get("archive_path") for d in config_docs]
        self.assertIn("dir1/config.txt", paths)
        self.assertIn("dir2/config.txt", paths)
        self.assertIn("dir3/config.txt", paths)

    def test_file_count_limit(self) -> None:
        archive_path = Path(self.temp_dir.name) / "many_files.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            for i in range(100):
                archive.writestr(f"file_{i}.txt", f"内容 {i}")
        
        limited_reader = ZipReader(max_files=50)
        with self.assertRaises(ValueError) as context:
            limited_reader._load_data(archive_path)
        self.assertIn("maximum file count", str(context.exception))

    def test_total_size_limit(self) -> None:
        archive_path = Path(self.temp_dir.name) / "large_total.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            for i in range(20):
                archive.writestr(f"file_{i}.txt", "x" * 1000)
        
        limited_reader = ZipReader(max_total_size=5000)
        with self.assertRaises(ValueError) as context:
            limited_reader._load_data(archive_path)
        self.assertIn("maximum total size", str(context.exception))

    def test_path_traversal_protection(self) -> None:
        archive_path = Path(self.temp_dir.name) / "traversal.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("../../../etc/passwd", "should be blocked")
            archive.writestr("./../../sensitive.txt", "should be blocked")
            archive.writestr("normal/file.txt", "正常文件")
            archive.writestr("../outside.txt", "应该被阻止")
        
        docs = self.reader._load_data(archive_path)
        
        for doc in docs:
            path = doc.metadata.get("archive_path", "")
            self.assertNotIn("..", path)

    def test_unsafe_members_are_skipped_not_normalized(self) -> None:
        archive_path = Path(self.temp_dir.name) / "unsafe_members.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("../outside.txt", "blocked")
            archive.writestr("/absolute.txt", "blocked")
            archive.writestr("safe/file.txt", "allowed")

        docs = self.reader._load_data(archive_path)

        self.assertEqual(len(docs), 1)
        self.assertEqual(docs[0].metadata["archive_path"], "safe/file.txt")
        self.assertEqual(docs[0].text, "allowed")

    def test_hidden_and_system_files(self) -> None:
        archive_path = Path(self.temp_dir.name) / "hidden.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr(".hidden", "隐藏文件")
            archive.writestr(".gitignore", "*.pyc\n__pycache__/")
            archive.writestr(".env", "SECRET_KEY=abc123")
            archive.writestr("normal.txt", "正常文件")
            archive.writestr("dir/.hidden_in_dir", "目录中的隐藏文件")
        
        docs = self.reader._load_data(archive_path)
        self.assertGreater(len(docs), 0)

    def test_unicode_content(self) -> None:
        archive_path = Path(self.temp_dir.name) / "unicode.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("chinese.txt", "这是中文内容：你好世界！")
            archive.writestr("japanese.txt", "日本語のテキスト：こんにちは")
            archive.writestr("korean.txt", "한국어 텍스트: 안녕하세요")
            archive.writestr("emoji.txt", "表情符号测试 😀 🎉 ✨ 🚀")
            archive.writestr("mixed.txt", "混合内容 Mixed Content مرحبا Привет")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 5)
        
        chinese_doc = [d for d in docs if "chinese.txt" in d.metadata.get("file_name", "")][0]
        self.assertIn("你好世界", chinese_doc.text)

    def test_various_compression_levels(self) -> None:
        content = "重复内容 " * 100
        
        for compression in [zipfile.ZIP_STORED, zipfile.ZIP_DEFLATED]:
            archive_path = Path(self.temp_dir.name) / f"compress_{compression}.zip"
            with zipfile.ZipFile(archive_path, "w", compression=compression) as archive:
                archive.writestr("data.txt", content)
            
            reader = ZipReader(max_compression_ratio=500)
            docs = reader._load_data(archive_path)
            self.assertEqual(len(docs), 1)
            self.assertIn("重复内容", docs[0].text)

    def test_empty_zip(self) -> None:
        archive_path = Path(self.temp_dir.name) / "empty.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            pass
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 0)

    def test_zip_with_only_directories(self) -> None:
        archive_path = Path(self.temp_dir.name) / "only_dirs.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("dir1/", "")
            archive.writestr("dir2/subdir/", "")
            archive.writestr("dir3/", "")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 0)

    def test_mixed_empty_and_content_files(self) -> None:
        archive_path = Path(self.temp_dir.name) / "mixed_empty.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("empty1.txt", "")
            archive.writestr("content.txt", "有内容")
            archive.writestr("empty2.md", "")
            archive.writestr("data.json", '{}')
            archive.writestr("empty3.py", "")
        
        docs = self.reader._load_data(archive_path)
        self.assertGreater(len(docs), 0)
        self.assertLessEqual(len(docs), 5)

    def test_very_long_filenames(self) -> None:
        archive_path = Path(self.temp_dir.name) / "long_names.zip"
        long_name = "a" * 200 + ".txt"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr(long_name, "内容")
            archive.writestr("dir/" + "b" * 150 + ".md", "# 标题")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 2)

    def test_multiple_nested_zips_same_level(self) -> None:
        archive_path = Path(self.temp_dir.name) / "multi_nested.zip"
        
        nested1 = io.BytesIO()
        with zipfile.ZipFile(nested1, "w") as z:
            z.writestr("data1.txt", "嵌套包1数据")
        
        nested2 = io.BytesIO()
        with zipfile.ZipFile(nested2, "w") as z:
            z.writestr("data2.txt", "嵌套包2数据")
        
        nested3 = io.BytesIO()
        with zipfile.ZipFile(nested3, "w") as z:
            z.writestr("data3.txt", "嵌套包3数据")
        
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("root.txt", "根文件")
            archive.writestr("archives/pack1.zip", nested1.getvalue())
            archive.writestr("archives/pack2.zip", nested2.getvalue())
            archive.writestr("archives/pack3.zip", nested3.getvalue())
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 4)
        
        nested_docs = [d for d in docs if ".zip/" in d.metadata.get("archive_path", "")]
        self.assertEqual(len(nested_docs), 3)

    def test_csv_parsing_in_zip(self) -> None:
        archive_path = Path(self.temp_dir.name) / "csv_test.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("data/sales.csv", "产品,数量,价格\n笔记本,100,5000\n鼠标,200,50")
            archive.writestr("data/users.csv", "用户名,邮箱\nzhangsan,zhang@test.com\nlisi,li@test.com")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 2)

    def test_json_and_yaml_in_nested_zip(self) -> None:
        archive_path = Path(self.temp_dir.name) / "config_archive.zip"
        
        nested = io.BytesIO()
        with zipfile.ZipFile(nested, "w") as z:
            z.writestr("app.json", '{"name": "app", "version": "2.0"}')
            z.writestr("db.yml", "host: localhost\nport: 3306")
        
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("main.json", '{"type": "main"}')
            archive.writestr("configs/nested.zip", nested.getvalue())
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 3)
        
        json_docs = [d for d in docs if d.metadata.get("file_name", "").endswith(".json")]
        self.assertEqual(len(json_docs), 2)

    def test_metadata_propagation_through_nesting(self) -> None:
        archive_path = Path(self.temp_dir.name) / "meta_nest.zip"
        
        nested = io.BytesIO()
        with zipfile.ZipFile(nested, "w") as z:
            z.writestr("inner.txt", "内部内容")
        
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("outer.txt", "外部内容")
            archive.writestr("nest/inner.zip", nested.getvalue())
        
        custom_meta = {
            "project": "测试项目",
            "version": "1.0",
            "author": "测试者"
        }
        
        docs = self.reader._load_data(archive_path, ext_info=custom_meta)
        
        for doc in docs:
            self.assertEqual(doc.metadata.get("project"), "测试项目")
            self.assertEqual(doc.metadata.get("version"), "1.0")
            self.assertEqual(doc.metadata.get("author"), "测试者")

    def test_archive_root_and_path_metadata(self) -> None:
        archive_path = Path(self.temp_dir.name) / "test_archive.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("level1/file.txt", "内容")
        
        docs = self.reader._load_data(archive_path)
        doc = docs[0]
        
        self.assertEqual(doc.metadata.get("archive_root"), "test_archive.zip")
        self.assertEqual(doc.metadata.get("archive_path"), "level1/file.txt")
        self.assertEqual(doc.metadata.get("file_name"), "file.txt")
        self.assertEqual(doc.metadata.get("archive_depth"), 0)

    def test_nested_archive_path_construction(self) -> None:
        archive_path = Path(self.temp_dir.name) / "path_test.zip"
        
        level2 = io.BytesIO()
        with zipfile.ZipFile(level2, "w") as z:
            z.writestr("deep/file.txt", "深层内容")
        
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("container/level2.zip", level2.getvalue())
        
        docs = self.reader._load_data(archive_path)
        doc = docs[0]
        
        expected_path = "container/level2.zip/deep/file.txt"
        self.assertEqual(doc.metadata.get("archive_path"), expected_path)
        self.assertEqual(doc.metadata.get("archive_depth"), 1)

    def test_large_number_of_small_files(self) -> None:
        archive_path = Path(self.temp_dir.name) / "many_small.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            for i in range(500):
                archive.writestr(f"files/batch_{i // 100}/file_{i}.txt", f"内容 {i}")
        
        docs = self.reader._load_data(archive_path)
        self.assertEqual(len(docs), 500)

    def test_whitespace_only_files(self) -> None:
        archive_path = Path(self.temp_dir.name) / "whitespace.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("spaces.txt", "   ")
            archive.writestr("tabs.txt", "\t\t\t")
            archive.writestr("newlines.txt", "\n\n\n")
            archive.writestr("mixed.txt", "  \n\t  \n  ")
            archive.writestr("normal.txt", "正常内容")
        
        docs = self.reader._load_data(archive_path)
        self.assertGreater(len(docs), 0)
        self.assertLessEqual(len(docs), 5)

    def test_binary_files_skipped(self) -> None:
        archive_path = Path(self.temp_dir.name) / "binary.zip"
        with zipfile.ZipFile(archive_path, "w") as archive:
            archive.writestr("image.png", bytes([0x89, 0x50, 0x4E, 0x47] + [0] * 100))
            archive.writestr("data.bin", bytes(range(256)))
            archive.writestr("text.txt", "文本内容")
        
        docs = self.reader._load_data(archive_path)
        
        text_docs = [d for d in docs if d.metadata.get("file_name") == "text.txt"]
        self.assertEqual(len(text_docs), 1)


if __name__ == "__main__":
    unittest.main()
