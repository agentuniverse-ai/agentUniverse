#!/usr/bin/env python3

"""PowerPoint presentation creation and inspection tool."""

# Validation failures are converted to structured tool responses at the public
# execute boundary, so bespoke exception subclasses would add no useful signal.
# ruff: noqa: TRY003, TRY004

import os
import tempfile
import zipfile
from typing import Any, ClassVar, cast

from agentuniverse.agent.action.tool.common_tool.file_path_utils import (
    resolve_safe_path,
)
from agentuniverse.agent.action.tool.tool import Tool


class PowerPointTool(Tool):
    """Create, append to, read, and inspect ``.pptx`` presentations.

    ``python-pptx`` is loaded only when an operation needs it. All file paths
    are confined to ``base_dir`` and writes use a temporary file followed by
    ``os.replace`` so a failed save cannot corrupt an existing presentation.

    Slide input is intentionally structured rather than accepting executable
    template code. Each slide may contain ``title``, ``subtitle``, ``bullets``,
    ``table``, ``notes``, and a supported ``layout`` name.
    """

    base_dir: str = "."
    max_read_bytes: int = 20 * 1024 * 1024
    max_write_bytes: int = 20 * 1024 * 1024
    max_uncompressed_bytes: int = 100 * 1024 * 1024
    max_archive_entries: int = 5_000
    max_slides: int = 100
    max_shapes_per_slide: int = 200
    max_text_chars: int = 50_000
    max_table_rows: int = 50
    max_table_columns: int = 20

    _LAYOUT_INDEX: ClassVar[dict[str, int]] = {
        "title": 0,
        "title_content": 1,
        "section": 2,
        "title_only": 5,
        "blank": 6,
    }
    _SLIDE_FIELDS: ClassVar[set[str]] = {
        "title",
        "subtitle",
        "bullets",
        "table",
        "notes",
        "layout",
    }
    _METADATA_FIELDS: ClassVar[set[str]] = {
        "title",
        "subject",
        "author",
        "keywords",
        "comments",
        "category",
    }
    _FALLBACK_TITLE_SHAPE_NAME: ClassVar[str] = "agentUniverse Title"

    def execute(
        self,
        mode: str,
        file_path: str,
        slides: list[dict[str, Any]] | None = None,
        overwrite: bool = False,
        template_path: str | None = None,
        metadata: dict[str, str] | None = None,
    ) -> dict[str, Any]:
        """Run a PowerPoint operation.

        Args:
            mode: One of ``create``, ``append``, ``read``, or ``info``.
            file_path: Presentation path, resolved underneath ``base_dir``.
            slides: Structured slide specifications for create/append.
            overwrite: Whether create mode may replace an existing file.
            template_path: Optional PPTX template used by create mode.
            metadata: Optional core presentation properties for create mode.

        Returns:
            A structured success or error dictionary.
        """
        try:
            self._validate_configuration()
            normalized_mode = self._normalize_mode(mode)
            safe_path = self._resolve_pptx_path(file_path, "file_path")

            if normalized_mode == "create":
                return self._create(
                    safe_path,
                    slides,
                    overwrite,
                    template_path,
                    metadata,
                )
            if normalized_mode == "append":
                return self._append(safe_path, slides)
            if normalized_mode == "read":
                return self._read(safe_path)
            return self._info(safe_path)
        except ImportError as exc:
            return self._error(
                file_path,
                "dependency_error",
                "python-pptx is required for PowerPoint operations. Install it with: pip install python-pptx",
                detail=str(exc),
            )
        except (TypeError, ValueError) as exc:
            return self._error(file_path, "validation_error", str(exc))
        except Exception as exc:  # library-specific parse/write failures
            return self._error(
                file_path,
                "operation_error",
                f"PowerPoint operation failed: {exc}",
            )

    @staticmethod
    def _error(
        file_path: Any,
        error_type: str,
        message: str,
        detail: str | None = None,
    ) -> dict[str, Any]:
        result = {
            "status": "error",
            "error_type": error_type,
            "error": message,
            "file_path": file_path,
        }
        if detail:
            result["detail"] = detail
        return result

    def _validate_configuration(self) -> None:
        limits = {
            "max_read_bytes": self.max_read_bytes,
            "max_write_bytes": self.max_write_bytes,
            "max_uncompressed_bytes": self.max_uncompressed_bytes,
            "max_archive_entries": self.max_archive_entries,
            "max_slides": self.max_slides,
            "max_shapes_per_slide": self.max_shapes_per_slide,
            "max_text_chars": self.max_text_chars,
            "max_table_rows": self.max_table_rows,
            "max_table_columns": self.max_table_columns,
        }
        for name, value in limits.items():
            if isinstance(value, bool) or not isinstance(value, int) or value <= 0:
                raise ValueError(f"{name} must be a positive integer")
        if not isinstance(self.base_dir, str) or not self.base_dir:
            raise ValueError("base_dir must be a non-empty string")

    @staticmethod
    def _normalize_mode(mode: str) -> str:
        if not isinstance(mode, str) or not mode.strip():
            raise ValueError("mode must be a non-empty string")
        normalized = mode.strip().lower()
        allowed = {"create", "append", "read", "info"}
        if normalized not in allowed:
            raise ValueError(f"invalid mode {mode!r}; expected one of: {', '.join(sorted(allowed))}")
        return normalized

    def _resolve_pptx_path(self, file_path: str, field_name: str) -> str:
        if not isinstance(file_path, str) or not file_path.strip():
            raise ValueError(f"{field_name} must be a non-empty string")
        if os.path.splitext(file_path)[1].lower() != ".pptx":
            raise ValueError(f"{field_name} must have a .pptx extension")
        return cast(str, resolve_safe_path(file_path, self.base_dir))

    def _ensure_readable(self, file_path: str, field_name: str = "file_path") -> None:
        if not os.path.isfile(file_path):
            raise ValueError(f"{field_name} does not exist: {file_path}")
        size = os.path.getsize(file_path)
        if size > self.max_read_bytes:
            raise ValueError(f"{field_name} size {size} bytes exceeds max_read_bytes ({self.max_read_bytes})")
        self._validate_archive_safety(file_path, field_name)

    def _validate_archive_safety(self, file_path: str, field_name: str = "file_path") -> None:
        """Bound the ZIP container before python-pptx expands it."""
        try:
            with zipfile.ZipFile(file_path) as archive:
                entries = archive.infolist()
                if len(entries) > self.max_archive_entries:
                    raise ValueError(
                        f"{field_name} contains {len(entries)} archive entries, "
                        f"exceeding max_archive_entries ({self.max_archive_entries})"
                    )
                uncompressed_size = sum(entry.file_size for entry in entries)
                if uncompressed_size > self.max_uncompressed_bytes:
                    raise ValueError(
                        f"{field_name} expands to {uncompressed_size} bytes, "
                        f"exceeding max_uncompressed_bytes ({self.max_uncompressed_bytes})"
                    )
        except zipfile.BadZipFile as exc:
            raise ValueError(f"{field_name} is not a valid PPTX ZIP archive") from exc

    @staticmethod
    def _load_pptx() -> tuple[Any, Any]:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:
            raise ImportError("No module named 'pptx'") from exc
        return Presentation, Inches

    def _create(
        self,
        file_path: str,
        slides: list[dict[str, Any]] | None,
        overwrite: bool,
        template_path: str | None,
        metadata: dict[str, str] | None,
    ) -> dict[str, Any]:
        if not isinstance(overwrite, bool):
            raise ValueError("overwrite must be a boolean")
        if os.path.exists(file_path) and not overwrite:
            raise ValueError(f"file already exists: {file_path}; set overwrite=true to replace it")

        validated_slides = self._validate_slides(slides, current_count=0)
        validated_metadata = self._validate_metadata(metadata)
        safe_template = None
        if template_path is not None:
            safe_template = self._resolve_pptx_path(template_path, "template_path")
            self._ensure_readable(safe_template, "template_path")

        Presentation, Inches = self._load_pptx()
        presentation = Presentation(safe_template) if safe_template else Presentation()
        if len(presentation.slides) + len(validated_slides) > self.max_slides:
            raise ValueError(f"template slides plus requested slides exceed max_slides ({self.max_slides})")
        self._apply_metadata(presentation, validated_metadata)
        for spec in validated_slides:
            self._add_slide(presentation, spec, Inches)

        self._atomic_save(presentation, file_path)
        return {
            "status": "success",
            "mode": "create",
            "file_path": file_path,
            "slides_added": len(validated_slides),
            "slide_count": len(presentation.slides),
            "file_size": os.path.getsize(file_path),
            "template_path": safe_template,
            "overwritten": overwrite,
        }

    def _append(
        self,
        file_path: str,
        slides: list[dict[str, Any]] | None,
    ) -> dict[str, Any]:
        self._ensure_readable(file_path)
        Presentation, Inches = self._load_pptx()
        presentation = Presentation(file_path)
        existing_count = len(presentation.slides)
        validated_slides = self._validate_slides(
            slides,
            current_count=existing_count,
        )
        for spec in validated_slides:
            self._add_slide(presentation, spec, Inches)
        self._atomic_save(presentation, file_path)
        return {
            "status": "success",
            "mode": "append",
            "file_path": file_path,
            "slides_added": len(validated_slides),
            "slide_count": len(presentation.slides),
            "file_size": os.path.getsize(file_path),
        }

    def _read(self, file_path: str) -> dict[str, Any]:  # noqa: C901
        self._ensure_readable(file_path)
        Presentation, _ = self._load_pptx()
        presentation = Presentation(file_path)
        self._validate_presentation_slide_count(presentation)

        # Read-side budgets mirror create/append: a compact XML presentation
        # can still expand into a huge Python structure if every shape/table/
        # row/cell is walked, so traversal stops the moment any budget is
        # exhausted instead of padding the result with empty strings.
        budget = _ReadBudget(self)
        slide_results: list[dict[str, Any]] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_shape = self._get_title_shape(slide)
            title = budget.consume_text(title_shape.text if title_shape is not None else "")[0]
            texts: list[str] = []
            tables: list[list[list[str]]] = []
            shape_count = 0
            for shape in slide.shapes:
                if shape_count >= self.max_shapes_per_slide:
                    budget.mark_truncated()
                    break
                if getattr(shape, "has_table", False):
                    shape_count += 1
                    tables.append(self._read_table(shape.table, budget))
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                if title_shape is not None and shape.shape_id == title_shape.shape_id:
                    continue
                shape_count += 1
                value = budget.consume_text(shape.text)[0]
                if value:
                    texts.append(value)
                if budget.chars_exhausted():
                    budget.mark_truncated()
                    break

            notes = ""
            if getattr(slide, "has_notes_slide", False):
                notes_frame = slide.notes_slide.notes_text_frame
                notes = budget.consume_text(notes_frame.text if notes_frame is not None else "")[0]

            slide_results.append(
                {
                    "slide_number": slide_number,
                    "title": title,
                    "texts": texts,
                    "tables": tables,
                    "notes": notes,
                }
            )
            if budget.chars_exhausted():
                break

        return {
            "status": "success",
            "mode": "read",
            "file_path": file_path,
            "slide_count": len(presentation.slides),
            "slides": slide_results,
            "truncated": budget.truncated,
            "max_text_chars": self.max_text_chars,
            "max_slides": self.max_slides,
            "max_shapes_per_slide": self.max_shapes_per_slide,
            "max_table_rows": self.max_table_rows,
            "max_table_columns": self.max_table_columns,
        }

    def _read_table(self, table: Any, budget: "_ReadBudget") -> list[list[str]]:
        rows: list[list[str]] = []
        for row_index, row in enumerate(table.rows):
            if row_index >= self.max_table_rows:
                budget.mark_truncated()
                break
            values: list[str] = []
            for column_index, cell in enumerate(row.cells):
                if column_index >= self.max_table_columns:
                    budget.mark_truncated()
                    break
                values.append(budget.consume_text(cell.text)[0])
            rows.append(values)
            if budget.chars_exhausted():
                budget.mark_truncated()
                break
        return rows

    def _info(self, file_path: str) -> dict[str, Any]:
        self._ensure_readable(file_path)
        Presentation, _ = self._load_pptx()
        presentation = Presentation(file_path)
        self._validate_presentation_slide_count(presentation)

        slide_info = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_shape = self._get_title_shape(slide)
            title = title_shape.text.strip() if title_shape is not None else ""
            slide_info.append(
                {
                    "slide_number": slide_number,
                    "title": title[:500],
                    "shape_count": len(slide.shapes),
                    "table_count": sum(1 for shape in slide.shapes if getattr(shape, "has_table", False)),
                    "has_notes": bool(getattr(slide, "has_notes_slide", False)),
                }
            )

        props = presentation.core_properties
        return {
            "status": "success",
            "mode": "info",
            "file_path": file_path,
            "file_size": os.path.getsize(file_path),
            "slide_count": len(presentation.slides),
            "slide_width": int(presentation.slide_width),
            "slide_height": int(presentation.slide_height),
            "metadata": {field: getattr(props, field, "") or "" for field in sorted(self._METADATA_FIELDS)},
            "slides": slide_info,
        }

    def _validate_presentation_slide_count(self, presentation: Any) -> None:
        count = len(presentation.slides)
        if count > self.max_slides:
            raise ValueError(f"presentation has {count} slides, exceeding max_slides ({self.max_slides})")

    def _validate_slides(  # noqa: C901
        self,
        slides: list[dict[str, Any]] | None,
        current_count: int,
    ) -> list[dict[str, Any]]:
        if not isinstance(slides, list) or not slides:
            raise ValueError("slides must be a non-empty list")
        if current_count + len(slides) > self.max_slides:
            raise ValueError(
                f"result would have {current_count + len(slides)} slides, exceeding max_slides ({self.max_slides})"
            )

        total_chars = 0
        validated = []
        for index, raw_spec in enumerate(slides, start=1):
            if not isinstance(raw_spec, dict):
                raise ValueError(f"slides[{index - 1}] must be an object")
            unknown = set(raw_spec) - self._SLIDE_FIELDS
            if unknown:
                raise ValueError(f"slides[{index - 1}] has unknown fields: {', '.join(sorted(unknown))}")
            spec = dict(raw_spec)
            for field in ("title", "subtitle", "notes"):
                value = spec.get(field, "")
                if not isinstance(value, str):
                    raise ValueError(f"slides[{index - 1}].{field} must be a string")
                spec[field] = value
                total_chars += len(value)

            layout = spec.get("layout", "auto")
            if not isinstance(layout, str):
                raise ValueError(f"slides[{index - 1}].layout must be a string")
            layout = layout.strip().lower()
            allowed_layouts = {"auto", *self._LAYOUT_INDEX}
            if layout not in allowed_layouts:
                raise ValueError(f"slides[{index - 1}].layout must be one of: {', '.join(sorted(allowed_layouts))}")
            spec["layout"] = layout

            bullets = spec.get("bullets", [])
            if not isinstance(bullets, list):
                raise ValueError(f"slides[{index - 1}].bullets must be a list")
            normalized_bullets = []
            for bullet_index, bullet in enumerate(bullets):
                if isinstance(bullet, str):
                    text, level = bullet, 0
                elif isinstance(bullet, dict):
                    unknown_bullet = set(bullet) - {"text", "level"}
                    if unknown_bullet:
                        raise ValueError(
                            f"slides[{index - 1}].bullets[{bullet_index}] has "
                            f"unknown fields: {', '.join(sorted(unknown_bullet))}"
                        )
                    raw_text = bullet.get("text")
                    level = bullet.get("level", 0)
                    if not isinstance(raw_text, str):
                        raise ValueError(f"slides[{index - 1}].bullets[{bullet_index}].text must be a string")
                    text = raw_text
                else:
                    raise ValueError(f"slides[{index - 1}].bullets[{bullet_index}] must be a string or object")
                if isinstance(level, bool) or not isinstance(level, int) or not 0 <= level <= 8:
                    raise ValueError(
                        f"slides[{index - 1}].bullets[{bullet_index}].level must be an integer from 0 to 8"
                    )
                normalized_bullets.append({"text": text, "level": level})
                total_chars += len(text)
            spec["bullets"] = normalized_bullets

            table = spec.get("table", [])
            if table is None:
                table = []
            if not isinstance(table, list):
                raise ValueError(f"slides[{index - 1}].table must be a list")
            if len(table) > self.max_table_rows:
                raise ValueError(f"slides[{index - 1}].table exceeds max_table_rows ({self.max_table_rows})")
            max_columns = 0
            normalized_table = []
            for row_index, row in enumerate(table):
                if not isinstance(row, list):
                    raise ValueError(f"slides[{index - 1}].table[{row_index}] must be a list")
                max_columns = max(max_columns, len(row))
                if max_columns > self.max_table_columns:
                    raise ValueError(f"slides[{index - 1}].table exceeds max_table_columns ({self.max_table_columns})")
                output_row = []
                for value in row:
                    if value is None:
                        value = ""
                    elif not isinstance(value, (str, int, float, bool)):
                        raise ValueError(f"slides[{index - 1}].table cells must be scalar values")
                    value = str(value)
                    output_row.append(value)
                    total_chars += len(value)
                normalized_table.append(output_row)
            if table and max_columns == 0:
                raise ValueError(f"slides[{index - 1}].table must contain cells")
            spec["table"] = normalized_table

            if not any((spec["title"], spec["subtitle"], spec["bullets"], spec["table"], spec["notes"])):
                raise ValueError(f"slides[{index - 1}] must contain content")
            validated.append(spec)

        if total_chars > self.max_text_chars:
            raise ValueError(
                f"slide content has {total_chars} characters, exceeding max_text_chars ({self.max_text_chars})"
            )
        return validated

    def _validate_metadata(
        self,
        metadata: dict[str, str] | None,
    ) -> dict[str, str]:
        if metadata is None:
            return {}
        if not isinstance(metadata, dict):
            raise ValueError("metadata must be an object")
        unknown = set(metadata) - self._METADATA_FIELDS
        if unknown:
            raise ValueError(f"metadata has unknown fields: {', '.join(sorted(unknown))}")
        result = {}
        for field, value in metadata.items():
            if not isinstance(value, str):
                raise ValueError(f"metadata.{field} must be a string")
            if len(value) > 2_000:
                raise ValueError(f"metadata.{field} exceeds 2000 characters")
            result[field] = value
        return result

    @staticmethod
    def _apply_metadata(presentation: Any, metadata: dict[str, str]) -> None:
        props = presentation.core_properties
        for field, value in metadata.items():
            setattr(props, field, value)

    def _add_slide(  # noqa: C901
        self,
        presentation: Any,
        spec: dict[str, Any],
        Inches: Any,
    ) -> None:
        layout_name = spec["layout"]
        if layout_name == "auto":
            layout_name = "title" if spec["subtitle"] and not spec["bullets"] and not spec["table"] else "title_content"
        requested_index = self._LAYOUT_INDEX[layout_name]
        layout_index = requested_index if requested_index < len(presentation.slide_layouts) else 0
        slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])

        title_shape = slide.shapes.title
        if spec["title"]:
            if title_shape is not None:
                title_shape.text = spec["title"]
            else:
                title_shape = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8.6), Inches(0.8))
                title_shape.name = self._FALLBACK_TITLE_SHAPE_NAME
                title_shape.text_frame.text = spec["title"]

        if spec["subtitle"]:
            subtitle_shape = self._find_body_placeholder(slide, title_shape)
            if layout_name in {"title", "section"} and subtitle_shape is not None:
                subtitle_shape.text_frame.text = spec["subtitle"]
            else:
                subtitle_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.55))
                subtitle_shape.text_frame.text = spec["subtitle"]

        if spec["bullets"]:
            body_shape = self._find_body_placeholder(slide, title_shape)
            if body_shape is None or (spec["subtitle"] and layout_name in {"title", "section"}):
                body_top = 1.9 if spec["subtitle"] else 1.45
                body_height = 2.0 if spec["table"] else 5.2
                body_shape = slide.shapes.add_textbox(Inches(0.9), Inches(body_top), Inches(8.2), Inches(body_height))
            self._set_bullets(body_shape.text_frame, spec["bullets"])

        if spec["table"]:
            rows = len(spec["table"])
            columns = max(len(row) for row in spec["table"])
            top = 4.0 if spec["bullets"] else (2.0 if spec["subtitle"] else 1.6)
            height = max(0.7, min(3.0, rows * 0.45))
            table_shape = slide.shapes.add_table(
                rows,
                columns,
                Inches(0.7),
                Inches(top),
                Inches(8.6),
                Inches(height),
            )
            table = table_shape.table
            for row_index, row in enumerate(spec["table"]):
                for column_index in range(columns):
                    table.cell(row_index, column_index).text = row[column_index] if column_index < len(row) else ""

        if spec["notes"]:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                notes_frame.text = spec["notes"]

    @staticmethod
    def _find_body_placeholder(slide: Any, title_shape: Any) -> Any | None:
        for shape in slide.placeholders:
            if not getattr(shape, "has_text_frame", False):
                continue
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            return shape
        return None

    def _get_title_shape(self, slide: Any) -> Any | None:
        """Return a native title placeholder or a title textbox created by this tool."""
        native_title = slide.shapes.title
        if native_title is not None:
            return native_title
        return next(
            (shape for shape in slide.shapes if shape.name == self._FALLBACK_TITLE_SHAPE_NAME),
            None,
        )

    @staticmethod
    def _set_bullets(text_frame: Any, bullets: list[dict[str, Any]]) -> None:
        text_frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet["text"]
            paragraph.level = bullet["level"]

    def _atomic_save(self, presentation: Any, file_path: str) -> None:
        directory = os.path.dirname(file_path)
        os.makedirs(directory, exist_ok=True)
        temporary_path = None
        try:
            with tempfile.NamedTemporaryFile(
                prefix=".powerpoint-",
                suffix=".pptx",
                dir=directory,
                delete=False,
            ) as temporary:
                temporary_path = temporary.name
            presentation.save(temporary_path)
            size = os.path.getsize(temporary_path)
            if size > self.max_write_bytes:
                raise ValueError(f"generated file size {size} bytes exceeds max_write_bytes ({self.max_write_bytes})")
            self._validate_archive_safety(temporary_path, "generated presentation")
            os.replace(temporary_path, file_path)
            temporary_path = None
        finally:
            if temporary_path and os.path.exists(temporary_path):
                os.unlink(temporary_path)

    @staticmethod
    def _bounded_text(text: Any, remaining: int) -> tuple[str, int, bool]:
        normalized = str(text or "").strip()
        if not normalized:
            return "", remaining, False
        if len(normalized) <= remaining:
            return normalized, remaining - len(normalized), False
        if remaining <= 0:
            return "", 0, True
        if remaining == 1:
            return "…", 0, True
        return normalized[: remaining - 1] + "…", 0, True


class _ReadBudget:
    """Tracks the char budget consumed across a presentation read.

    A crafted PPTX can pack in many slides, shapes, tables, rows, and cells.
    This accumulator centralises the text-char budget so the read stops the
    moment it is exhausted, instead of continuing to walk every shape and
    padding the result with empty strings. Slide/table/row/column caps are
    enforced inline against the tool's configured limits; this budget owns the
    shared text budget that spans all of them. Mirrors the WordDocumentTool
    read contract.
    """

    __slots__ = ("remaining_chars", "truncated")

    def __init__(self, tool: "PowerPointTool") -> None:
        self.remaining_chars = tool.max_text_chars
        self.truncated = False

    def chars_exhausted(self) -> bool:
        return self.remaining_chars <= 0

    def mark_truncated(self) -> None:
        self.truncated = True

    def consume_text(self, value: Any) -> tuple[str, bool]:
        """Return (bounded_text, was_cut). Empty results do not consume budget."""
        normalized = str(value or "").strip()
        if not normalized:
            return "", False
        if len(normalized) <= self.remaining_chars:
            self.remaining_chars -= len(normalized)
            return normalized, False
        if self.remaining_chars <= 0:
            return "", True
        if self.remaining_chars == 1:
            self.remaining_chars = 0
            return "…", True
        kept = self.remaining_chars
        self.remaining_chars = 0
        return normalized[: max(0, kept - 1)] + "…", True
